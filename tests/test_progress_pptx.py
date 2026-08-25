"""Чтение PPTX-плана: текст фигур, отделение заголовка слайда, группировка в строки.

Презентация строится в памяти через python-pptx — как и предполагает формат
входа (карточки задач, таймлайн), без OCR/скриншотов.
"""

import tempfile
import unittest
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt

from pko.progress.pptx_reader import SlideShape, cluster_rows, read_deck


def _add_box(slide, left, top, width, height, title, body=""):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].text = title
    if body:
        p2 = tf.add_paragraph()
        p2.text = body
    return box


def build_sample_deck(path: Path) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    s1 = prs.slides.add_slide(blank)
    _add_box(s1, 1, 2.5, 11, 1.5, "Roadmap: редизайн платёжного сервиса", "Команда Payments")

    s2 = prs.slides.add_slide(blank)
    _add_box(s2, 0.5, 0.3, 6, 0.6, "Задачи спринта")
    cards = [
        ("Авторизация пользователей", "OAuth2 + JWT"),
        ("API платежей", "REST эндпоинты"),
        ("Уведомления", "Email + push"),
        ("Логирование", "Централизованный лог"),
        ("Ретраи операций", "Повтор через очередь"),
        ("Админ-дашборд", "Мониторинг платежей"),
    ]
    x0, y0, cw, ch = 0.5, 1.2, 4.0, 2.2
    for i, (title, body) in enumerate(cards):
        r, c = divmod(i, 3)
        _add_box(s2, x0 + c * (cw + 0.3), y0 + r * (ch + 0.3), cw, ch, title, body)

    s3 = prs.slides.add_slide(blank)
    _add_box(s3, 0.5, 0.3, 6, 0.6, "Таймлайн разработки")
    stages = [
        ("Этап 1: MVP", "июль"),
        ("Этап 2: Уведомления", "август"),
        ("Этап 3: Надёжность", "сентябрь"),
        ("Этап 4: Аналитика", "октябрь"),
    ]
    tx0, ty0, tw, th = 0.5, 2.5, 2.9, 1.6
    for i, (title, body) in enumerate(stages):
        _add_box(s3, tx0 + i * (tw + 0.2), ty0, tw, th, title, body)

    s4 = prs.slides.add_slide(blank)  # пустой слайд — потенциальный скриншот

    prs.save(str(path))


class PptxReaderTest(unittest.TestCase):
    @classmethod
    def setUpClass(cls):
        cls.tmp = tempfile.TemporaryDirectory()
        cls.path = Path(cls.tmp.name) / "plan.pptx"
        build_sample_deck(cls.path)
        cls.slides = read_deck(cls.path)

    @classmethod
    def tearDownClass(cls):
        cls.tmp.cleanup()

    def test_reads_all_slides_in_order(self):
        self.assertEqual([s.number for s in self.slides], [1, 2, 3, 4])

    def test_title_slide_has_no_split_heading(self):
        # Одна фигура — разделять не на что, весь текст остаётся в shapes.
        title_slide = self.slides[0]
        self.assertIsNone(title_slide.heading)
        self.assertEqual(len(title_slide.shapes), 1)
        self.assertIn("Roadmap", title_slide.shapes[0].text)

    def test_card_grid_heading_is_separated_from_cards(self):
        cards_slide = self.slides[1]
        self.assertEqual(cards_slide.heading, "Задачи спринта")
        self.assertEqual(len(cards_slide.shapes), 6)
        # 2 строки по 3 карточки — не один плоский список из шести.
        rows = cards_slide.rows
        self.assertEqual([len(r) for r in rows], [3, 3])
        self.assertEqual([s.text.splitlines()[0] for s in rows[0]],
                         ["Авторизация пользователей", "API платежей", "Уведомления"])

    def test_timeline_heading_does_not_pollute_row_clustering(self):
        # Регрессия на находку спайка Фазы 0: без отделения заголовка строка с
        # таймлайном ошибочно дробилась заголовком слайда сверху.
        timeline_slide = self.slides[2]
        self.assertEqual(timeline_slide.heading, "Таймлайн разработки")
        self.assertEqual(len(timeline_slide.shapes), 4)
        rows = timeline_slide.rows
        self.assertEqual(len(rows), 1, msg=f"ожидалась одна строка из 4 этапов: {rows}")
        self.assertEqual(len(rows[0]), 4)

    def test_empty_slide_is_flagged_not_silently_dropped(self):
        empty_slide = self.slides[3]
        self.assertTrue(empty_slide.is_empty)
        self.assertIsNone(empty_slide.heading)
        self.assertEqual(empty_slide.shapes, [])


class GroupedShapesTest(unittest.TestCase):
    """Найденный вручную баг: `slide.shapes` не спускается в группу сама по
    себе — без рекурсии текст сгруппированных фигур был не виден вообще, и
    слайд с реальным содержимым выглядел неотличимо от пустого."""

    def test_text_inside_a_group_is_still_read(self):
        with tempfile.TemporaryDirectory() as raw:
            path = Path(raw) / "grouped.pptx"
            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            b1 = _add_box(slide, 1, 1, 2, 1, "Приём событий", "REST API")
            b2 = _add_box(slide, 4, 1, 2, 1, "Email-канал", "Отправка почты")
            slide.shapes.add_group_shape([b1, b2])
            prs.save(str(path))

            slides = read_deck(path)
            self.assertEqual(len(slides), 1)
            self.assertFalse(slides[0].is_empty)
            texts = {s.text.splitlines()[0] for s in slides[0].shapes}
            self.assertEqual(texts, {"Приём событий", "Email-канал"})


class ClusterRowsTest(unittest.TestCase):
    """`cluster_rows` — группировка по близости `top`, отдельно от чтения PPTX."""

    def test_two_sections_on_one_slide_stay_in_separate_rows(self):
        # Ровно сценарий "всё стилизовано на одном слайде": карточки задач и
        # этапы таймлайна на разной высоте не должны слипаться в одну строку.
        cards = [SlideShape(f"Карточка {i}", left=i * 3.0, top=1.1, width=2.5, height=1.0)
                for i in range(3)]
        stages = [SlideShape(f"Этап {i}", left=i * 3.0, top=2.8, width=2.5, height=1.0)
                 for i in range(3)]
        rows = cluster_rows(cards + stages)
        self.assertEqual(len(rows), 2)
        self.assertEqual({s.text for s in rows[0]}, {s.text for s in cards})
        self.assertEqual({s.text for s in rows[1]}, {s.text for s in stages})

    def test_single_shape_is_a_row_of_one(self):
        self.assertEqual(cluster_rows([]), [])
        one = [SlideShape("x", 0, 0, 1, 1)]
        self.assertEqual(cluster_rows(one), [one])

    def test_row_items_are_ordered_left_to_right(self):
        shapes = [
            SlideShape("C", left=6, top=1, width=1, height=1),
            SlideShape("A", left=0, top=1, width=1, height=1),
            SlideShape("B", left=3, top=1, width=1, height=1),
        ]
        rows = cluster_rows(shapes)
        self.assertEqual([s.text for s in rows[0]], ["A", "B", "C"])


if __name__ == "__main__":
    unittest.main()
