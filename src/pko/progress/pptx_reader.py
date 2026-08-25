"""Чтение PPTX с планом: текст фигур по слайдам вместе с их расположением.

Формат входа предполагается текстовым (карточки задач, таймлайн-схема как
нативные фигуры PowerPoint) — см. Фазу 0 плана. OCR/vision здесь нет: слайд без
единой текстовой фигуры — сигнал «возможно, это картинка», а не пустой слайд,
и он помечается отдельно (`Slide.is_empty`), а не проваливается молча.
"""

from __future__ import annotations

from dataclasses import dataclass, field
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu


@dataclass(frozen=True)
class SlideShape:
    """Текст одной фигуры и её положение на слайде, в дюймах."""

    text: str
    left: float
    top: float
    width: float
    height: float


@dataclass(frozen=True)
class LayoutGuess:
    kind: str  # "SINGLE" | "TIMELINE" | "CARD_GRID" | "LIST"
    reason: str


@dataclass(frozen=True)
class Slide:
    number: int
    heading: str | None
    shapes: list[SlideShape] = field(default_factory=list)  # без фигуры-заголовка

    @property
    def is_empty(self) -> bool:
        return not self.heading and not self.shapes

    @property
    def layout(self) -> LayoutGuess:
        return guess_layout(self.shapes)


def read_deck(path: str | Path) -> list[Slide]:
    """Прочитать презентацию. Ничего не пишет и не мутирует исходный файл."""
    prs = Presentation(str(path))
    slides: list[Slide] = []
    for i, slide in enumerate(prs.slides, start=1):
        raw_shapes = _text_shapes(slide)
        heading, rest = _split_heading(raw_shapes)
        slides.append(Slide(number=i, heading=heading, shapes=rest))
    return slides


def _text_shapes(slide) -> list[SlideShape]:
    out: list[SlideShape] = []
    for shp in slide.shapes:
        if not getattr(shp, "has_text_frame", False):
            continue
        text = shp.text_frame.text.strip()
        if not text:
            continue
        out.append(SlideShape(
            text=text,
            left=_inches(shp.left),
            top=_inches(shp.top),
            width=_inches(shp.width),
            height=_inches(shp.height),
        ))
    return out


def _inches(emu: int | None) -> float:
    return round(Emu(emu).inches, 2) if emu is not None else 0.0


# Подпись слайда обычно короче одной строки и не переносится — этого достаточно,
# чтобы отличить «Таймлайн разработки» от текста внутри карточки задачи.
_HEADING_MAX_CHARS = 80


def _split_heading(shapes: list[SlideShape]) -> tuple[str | None, list[SlideShape]]:
    """Отделить подпись слайда от контентных фигур по позиции и форме.

    Заголовок — самая верхняя фигура, короткая и целиком выше остальных.
    Без этого шага таймлайн (все фигуры на одной высоте) вместе с подписью над
    ним (на другой высоте) выглядит как «разброс и по top, и по left» и
    ошибочно читается как сетка карточек — это показал спайк на синтетической
    презентации: слайд с таймлайном без разбора заголовка классифицировался
    как CARD_GRID.
    """
    if len(shapes) < 2:
        return None, list(shapes)
    ordered = sorted(shapes, key=lambda s: s.top)
    head, rest = ordered[0], ordered[1:]
    is_short = "\n" not in head.text and len(head.text) <= _HEADING_MAX_CHARS
    is_above_rest = all(head.top + head.height <= s.top + 0.05 for s in rest)
    if is_short and is_above_rest:
        return head.text, rest
    return None, list(shapes)


def guess_layout(shapes: list[SlideShape]) -> LayoutGuess:
    """Грубая эвристика раскладки — подсказка для LLM-роли planner, не факт.

    Точность здесь не обязана быть высокой: раскладка идёт в промпт как
    ориентир («вероятно, таймлайн»), а не как единственный источник структуры.
    """
    if len(shapes) <= 1:
        return LayoutGuess("SINGLE", "одна фигура или меньше")
    tops = sorted(s.top for s in shapes)
    lefts = sorted(s.left for s in shapes)
    top_spread = tops[-1] - tops[0]
    left_spread = lefts[-1] - lefts[0]
    if top_spread < 0.5 and left_spread > 1.0:
        return LayoutGuess("TIMELINE", "фигуры на одной высоте, разнесены по горизонтали")
    if top_spread > 1.0 and left_spread > 1.0:
        return LayoutGuess("CARD_GRID", "фигуры варьируются и по высоте, и по горизонтали")
    return LayoutGuess("LIST", "фигуры идут одна под другой")
