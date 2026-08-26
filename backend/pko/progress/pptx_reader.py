"""Чтение PPTX с планом: текст фигур по слайдам, сгруппированный в строки.

Формат входа предполагается текстовым (карточки задач, таймлайн-схема как
нативные фигуры PowerPoint) — см. Фазу 0 плана. OCR/vision здесь нет: слайд без
единой текстовой фигуры — сигнал «возможно, это картинка», а не пустой слайд,
и он помечается отдельно (`Slide.is_empty`), а не проваливается молча.

Группированные фигуры PowerPoint («оформить стилем» на практике часто значит
сгруппировать карточку целиком) разворачиваются рекурсивно — без этого их текст
не виден вообще: `slide.shapes` не спускается внутрь группы сам по себе, и
слайд с реальным содержимым выглядел бы неотличимо от пустого.
"""

from __future__ import annotations

from dataclasses import dataclass, field
from pathlib import Path
from typing import Iterable

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Emu


@dataclass(frozen=True)
class SlideShape:
    """Текст одной фигуры и её положение на слайде, в дюймах."""

    text: str
    left: float
    top: float
    width: float
    height: float


@dataclass(frozen=True)
class Slide:
    number: int
    heading: str | None
    shapes: list[SlideShape] = field(default_factory=list)  # без фигуры-заголовка

    @property
    def is_empty(self) -> bool:
        return not self.heading and not self.shapes

    @property
    def rows(self) -> list[list[SlideShape]]:
        return cluster_rows(self.shapes)


def read_deck(path: str | Path) -> list[Slide]:
    """Прочитать презентацию. Ничего не пишет и не мутирует исходный файл."""
    prs = Presentation(str(path))
    slides: list[Slide] = []
    for i, slide in enumerate(prs.slides, start=1):
        raw_shapes = _text_shapes(slide.shapes)
        heading, rest = _split_heading(raw_shapes)
        slides.append(Slide(number=i, heading=heading, shapes=rest))
    return slides


def _text_shapes(shapes: Iterable) -> list[SlideShape]:
    out: list[SlideShape] = []
    for shp in shapes:
        if getattr(shp, "shape_type", None) == MSO_SHAPE_TYPE.GROUP:
            # Координаты дочерних фигур python-pptx уже отдаёт в системе
            # слайда, не группы — пересчитывать смещение не нужно.
            out.extend(_text_shapes(shp.shapes))
            continue
        if not getattr(shp, "has_text_frame", False):
            continue
        text = shp.text_frame.text.strip()
        if not text:
            continue
        out.append(SlideShape(
            text=text,
            left=_inches(shp.left),
            top=_inches(shp.top),
            width=_inches(shp.width),
            height=_inches(shp.height),
        ))
    return out


def _inches(emu: int | None) -> float:
    return round(Emu(emu).inches, 2) if emu is not None else 0.0


# Подпись слайда обычно короче одной строки и не переносится — этого достаточно,
# чтобы отличить «Таймлайн разработки» от текста внутри карточки задачи.
_HEADING_MAX_CHARS = 80


def _split_heading(shapes: list[SlideShape]) -> tuple[str | None, list[SlideShape]]:
    """Отделить подпись слайда от контентных фигур по позиции и форме.

    Заголовок — самая верхняя фигура, короткая и целиком выше остальных.
    Без этого шага строка вычисляется вместе с подписью над ней, и она
    засчитывается как ещё одна «строка» из одной фигуры — расклад по строкам
    остаётся верным, но с лишней однострочной группой сверху.
    """
    if len(shapes) < 2:
        return None, list(shapes)
    ordered = sorted(shapes, key=lambda s: s.top)
    head, rest = ordered[0], ordered[1:]
    is_short = "\n" not in head.text and len(head.text) <= _HEADING_MAX_CHARS
    is_above_rest = all(head.top + head.height <= s.top + 0.05 for s in rest)
    if is_short and is_above_rest:
        return head.text, rest
    return None, list(shapes)


# Насколько близко по вертикали должны стоять фигуры, чтобы считаться одной
# строкой (ряд карточек, ряд этапов таймлайна). Больше — уже новая строка.
_ROW_TOP_TOLERANCE = 0.4


def cluster_rows(shapes: list[SlideShape]) -> list[list[SlideShape]]:
    """Сгруппировать фигуры по близости `top` — визуальные строки на слайде.

    Это группировка по расположению, а не по смыслу: код не решает, ряд
    карточек это или ряд этапов таймлайна — только то, что эти фигуры стоят
    на одной высоте и, скорее всего, составляют одну логическую строку.
    Смысл строки остаётся за LLM, у которого есть сам текст.

    Без этого шага несколько разных разделов на одном слайде (например,
    карточки задач и таймлайн, оформленные вместе без разделяющего
    заголовка) сливались в один плоский список фигур, и заголовок
    «поехавшей» строки было неоткуда взять — теперь строки видны отдельно
    друг от друга даже без заголовка между ними.
    """
    if not shapes:
        return []
    ordered = sorted(shapes, key=lambda s: (s.top, s.left))
    rows: list[list[SlideShape]] = [[ordered[0]]]
    for shp in ordered[1:]:
        if abs(shp.top - rows[-1][-1].top) <= _ROW_TOP_TOLERANCE:
            rows[-1].append(shp)
        else:
            rows.append([shp])
    for row in rows:
        row.sort(key=lambda s: s.left)
    return rows


def render_slide(slide: Slide) -> str:
    """Слайд как читаемый текст: заголовок, дальше — строки фигур по порядку.

    Не JSON: связный текст с явной структурой строк надёжнее для LLM, чем
    нагромождение чисел-координат — модель не vision, пространственное
    рассуждение по `left`/`top` для неё не сильная сторона. Строки уже
    сгруппированы кодом (`cluster_rows`, по близости `top`), а какая это по
    смыслу строка (задачи или этапы), решает модель по самому тексту.
    """
    header = f"Слайд {slide.number}"
    if slide.heading:
        header += f": {slide.heading}"
    lines = [header]

    rows = cluster_rows(slide.shapes)
    if not rows:
        lines.append("  (текстовых фигур на слайде нет)")
        return "\n".join(lines)

    for row_no, row in enumerate(rows, start=1):
        lines.append(f"  Строка {row_no} (фигур: {len(row)}):")
        for item_no, shape in enumerate(row, start=1):
            title, _, body = shape.text.partition("\n")
            lines.append(f"    {item_no}. {title.strip()}")
            if body.strip():
                lines.append(f"       {body.strip()}")
    return "\n".join(lines)
